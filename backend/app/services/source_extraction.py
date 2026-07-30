"""Format-aware extraction into chunks with exact, structured source locations."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import csv
from time import monotonic

from app.config import settings
from app.services.chunking import chunk_text
from app.services.document_loader import DocumentParseError, extract_text
from app.services.image_processor import IMAGE_EXTENSIONS, chunk_image_text


@dataclass(frozen=True)
class SourceChunk:
    text: str
    source_type: str
    location: dict[str, object]


def validate_source_location(source_type: str, location: dict[str, object]) -> None:
    """Reject chunks that cannot produce a precise citation for their source type."""
    required: dict[str, tuple[str, ...]] = {
        "pdf": ("page_start", "page_end"),
        "powerpoint": ("slide_start", "slide_end", "content_type"),
        "excel": (
            "sheet_name", "row_start", "row_end", "column_start",
            "column_end", "cell_range", "hidden_sheet",
        ),
        "csv": ("row_start", "row_end"),
        "text": ("line_start", "line_end"),
    }
    missing = [
        key for key in required.get(source_type, ())
        if key not in location or location[key] is None
    ]
    if missing:
        raise DocumentParseError(
            f"Invalid {source_type} source location; missing {', '.join(missing)}."
        )
    for start, end in (
        ("page_start", "page_end"),
        ("slide_start", "slide_end"),
        ("row_start", "row_end"),
        ("line_start", "line_end"),
    ):
        if start in location and end in location:
            if int(location[start]) < 1 or int(location[end]) < int(location[start]):
                raise DocumentParseError(
                    f"Invalid {source_type} source range: {start}/{end}."
                )


def _split(text: str, source_type: str, location: dict[str, object]) -> list[SourceChunk]:
    return [
        SourceChunk(value, source_type, {**location, "part": index})
        for index, value in enumerate(chunk_text(text), start=1)
    ]


def _pdf(path: Path) -> list[SourceChunk]:
    try:
        from pypdf import PdfReader

        result: list[SourceChunk] = []
        pages = PdfReader(str(path)).pages
        if len(pages) > settings.max_pdf_pages:
            raise DocumentParseError("PDF exceeds the configured page limit.")
        for page_number, page in enumerate(pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                for block_index, value in enumerate(chunk_text(text), start=1):
                    result.append(SourceChunk(value, "pdf", {
                        "page_start": page_number,
                        "page_end": page_number,
                        "block_ids": [f"p{page_number}-b{block_index}"],
                        "bounding_boxes": [],
                        "part": block_index,
                    }))
        return result
    except Exception as error:
        raise DocumentParseError("The PDF file could not be read.") from error


def _pptx(path: Path) -> list[SourceChunk]:
    try:
        from pptx import Presentation

        result: list[SourceChunk] = []
        slides = Presentation(str(path)).slides
        if len(slides) > settings.max_powerpoint_slides:
            raise DocumentParseError("Presentation exceeds the configured slide limit.")
        for slide_number, slide in enumerate(slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                shape_name = str(getattr(shape, "name", f"Shape {shape_index}"))
                shape_id = f"shape-{shape_index}"
                if shape == getattr(slide.shapes, "title", None):
                    shape_type = "title"
                elif getattr(shape, "has_table", False):
                    shape_type = "table"
                elif getattr(shape, "has_chart", False):
                    shape_type = "chart"
                else:
                    shape_type = "text_box"
                common = {
                    "slide_start": slide_number,
                    "slide_end": slide_number,
                    "slide_number": slide_number,
                    "shape_ids": [shape_id],
                    "shape_types": [shape_type],
                    "shape_name": shape_name,
                    "shape_index": shape_index,
                    "speaker_notes_included": False,
                }
                if getattr(shape, "has_table", False):
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        text = "\t".join(cell.text.strip() for cell in row.cells).strip()
                        if text:
                            result.append(SourceChunk(text, "powerpoint", {
                                **common, "content_type": "table",
                                "row_start": row_index, "row_end": row_index,
                            }))
                elif getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip() for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text:
                        result.extend(_split(text, "powerpoint", {
                            **common, "content_type": (
                                "chart_label" if shape_type == "chart" else "slide_text"
                            ),
                        }))
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    result.extend(_split(notes, "powerpoint", {
                        "slide_start": slide_number,
                        "slide_end": slide_number,
                        "slide_number": slide_number,
                        "shape_ids": [],
                        "shape_types": ["speaker_notes"],
                        "speaker_notes_included": True,
                        "content_type": "speaker_notes",
                    }))
            except (AttributeError, ValueError):
                pass
        return result
    except Exception as error:
        raise DocumentParseError("The PowerPoint presentation could not be read.") from error


def _xlsx(
    path: Path,
    include_hidden: bool,
    include_very_hidden: bool,
) -> list[SourceChunk]:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter, range_boundaries

        workbook = load_workbook(path, read_only=False, data_only=False)
        values_workbook = load_workbook(path, read_only=False, data_only=True)
        result: list[SourceChunk] = []
        try:
            if len(workbook.worksheets) > settings.max_workbook_sheets:
                raise DocumentParseError("Workbook exceeds the configured sheet limit.")
            total_rows = sum(sheet.max_row for sheet in workbook.worksheets)
            if total_rows > settings.max_workbook_rows:
                raise DocumentParseError("Workbook exceeds the configured row limit.")
            for worksheet in workbook.worksheets:
                values_sheet = values_workbook[worksheet.title]
                hidden = worksheet.sheet_state != "visible"
                if (
                    worksheet.sheet_state == "hidden" and not include_hidden
                    or worksheet.sheet_state == "veryHidden"
                    and not include_very_hidden
                ):
                    continue
                populated_rows = [
                    row for row in worksheet.iter_rows()
                    if any(cell.value is not None for cell in row)
                ]
                header_row_number = (
                    next(
                        cell.row
                        for row in populated_rows
                        for cell in row
                        if cell.value is not None
                    )
                    if populated_rows else None
                )
                header_context = [
                    str(cell.value)
                    for cell in (
                        worksheet[header_row_number] if header_row_number else []
                    )
                    if cell.value is not None
                ]
                tables: list[tuple[str, tuple[int, int, int, int]]] = [
                    (table.name, range_boundaries(table.ref))
                    for table in worksheet.tables.values()
                ]
                merged_ranges = [str(value) for value in worksheet.merged_cells.ranges]
                for row in populated_rows:
                    populated = [cell for cell in row if cell.value is not None]
                    if not populated:
                        continue
                    min_column = min(cell.column for cell in populated)
                    max_column = max(cell.column for cell in populated)
                    row_number = populated[0].row
                    cell_range = (
                        f"{get_column_letter(min_column)}{row_number}:"
                        f"{get_column_letter(max_column)}{row_number}"
                    )
                    formulas = {
                        cell.coordinate: str(cell.value)
                        for cell in populated
                        if isinstance(cell.value, str) and cell.value.startswith("=")
                    }
                    values = [
                        (
                            f"{cell.coordinate}: {cell.value}"
                            + (
                                f" (cached value: {values_sheet[cell.coordinate].value})"
                                if cell.coordinate in formulas
                                and values_sheet[cell.coordinate].value is not None
                                else ""
                            )
                        )
                        for cell in populated
                    ]
                    table_name = next(
                        (
                            name
                            for name, (table_min_col, table_min_row, table_max_col, table_max_row)
                            in tables
                            if table_min_row <= row_number <= table_max_row
                            and min_column >= table_min_col
                            and max_column <= table_max_col
                        ),
                        None,
                    )
                    result.append(SourceChunk("\t".join(values), "excel", {
                        "sheet_name": worksheet.title,
                        "hidden_sheet": hidden,
                        "sheet_hidden": hidden,
                        "row_start": row_number,
                        "row_end": row_number,
                        "column_start": get_column_letter(min_column),
                        "column_end": get_column_letter(max_column),
                        "cell_range": cell_range,
                        "table_name": table_name,
                        "header_rows": (
                            [header_row_number] if header_row_number else []
                        ),
                        "header_context": header_context,
                        "merged_ranges": merged_ranges,
                        "formulas": formulas,
                    }))
        finally:
            workbook.close()
            values_workbook.close()
        return result
    except Exception as error:
        raise DocumentParseError("The Excel workbook could not be read.") from error


def _xls(
    path: Path,
    include_hidden: bool,
    include_very_hidden: bool,
) -> list[SourceChunk]:
    try:
        import xlrd
        from openpyxl.utils import get_column_letter

        workbook = xlrd.open_workbook(path, on_demand=True)
        result: list[SourceChunk] = []
        try:
            for sheet in workbook.sheets():
                hidden = bool(getattr(sheet, "visibility", 0))
                visibility = int(getattr(sheet, "visibility", 0) or 0)
                if (
                    visibility == 1 and not include_hidden
                    or visibility >= 2 and not include_very_hidden
                ):
                    continue
                header_row = next(
                    (
                        row_index + 1
                        for row_index in range(sheet.nrows)
                        if any(sheet.cell_value(row_index, column) not in ("", None)
                               for column in range(sheet.ncols))
                    ),
                    None,
                )
                header_context = (
                    [
                        str(sheet.cell_value(header_row - 1, column))
                        for column in range(sheet.ncols)
                        if sheet.cell_value(header_row - 1, column) not in ("", None)
                    ]
                    if header_row else []
                )
                for row_index in range(sheet.nrows):
                    populated = [
                        column for column in range(sheet.ncols)
                        if sheet.cell_value(row_index, column) not in ("", None)
                    ]
                    if not populated:
                        continue
                    row_number = row_index + 1
                    column_start = get_column_letter(min(populated) + 1)
                    column_end = get_column_letter(max(populated) + 1)
                    result.append(SourceChunk(
                        "\t".join(
                            f"{get_column_letter(column + 1)}{row_number}: "
                            f"{sheet.cell_value(row_index, column)}"
                            for column in populated
                        ),
                        "excel",
                        {
                            "sheet_name": sheet.name,
                            "hidden_sheet": hidden,
                            "row_start": row_number,
                            "row_end": row_number,
                            "column_start": column_start,
                            "column_end": column_end,
                            "cell_range": (
                                f"{column_start}{row_number}:"
                                f"{column_end}{row_number}"
                            ),
                            "table_name": None,
                            "header_rows": [header_row] if header_row else [],
                            "header_context": header_context,
                            "merged_ranges": [],
                            "formulas": {},
                        },
                    ))
        finally:
            workbook.release_resources()
        return result
    except Exception as error:
        raise DocumentParseError("The Excel workbook could not be read.") from error


def _docx(path: Path) -> list[SourceChunk]:
    try:
        from docx import Document

        document = Document(str(path))
        result: list[SourceChunk] = []
        section_index = 1
        for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if not text:
                continue
            if paragraph._p.xpath("./w:pPr/w:sectPr"):
                section_index += 1
            result.extend(_split(text, "word", {
                "section_number": section_index,
                "paragraph_start": paragraph_index,
                "paragraph_end": paragraph_index,
                "style": paragraph.style.name if paragraph.style else None,
            }))
        for table_index, table in enumerate(document.tables, start=1):
            for row_index, row in enumerate(table.rows, start=1):
                text = "\t".join(cell.text.strip() for cell in row.cells).strip()
                if text:
                    result.append(SourceChunk(text, "word", {
                        "table_number": table_index,
                        "row_start": row_index,
                        "row_end": row_index,
                    }))
        return result
    except Exception as error:
        raise DocumentParseError("The DOCX file could not be read.") from error


def _csv(path: Path) -> list[SourceChunk]:
    try:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
            sample = handle.read(8192)
            handle.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel
            return [
                SourceChunk(
                    "\t".join(str(value).strip() for value in row),
                    "csv",
                    {"row_start": row_number, "row_end": row_number},
                )
                for row_number, row in enumerate(csv.reader(handle, dialect), start=1)
                if any(str(value).strip() for value in row)
            ]
    except Exception as error:
        raise DocumentParseError("The CSV file could not be read.") from error


def _text(path: Path) -> list[SourceChunk]:
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    result: list[SourceChunk] = []
    for line_number, value in enumerate(lines, start=1):
        if value.strip():
            result.append(SourceChunk(
                value.strip(), "text",
                {"line_start": line_number, "line_end": line_number},
            ))
    return result


def extract_source_chunks(
    path: Path,
    *,
    include_hidden: bool = True,
    include_very_hidden: bool = False,
) -> list[SourceChunk]:
    started_at = monotonic()
    extension = path.suffix.lower()
    if extension == ".pdf":
        chunks = _pdf(path)
    elif extension == ".pptx":
        chunks = _pptx(path)
    elif extension == ".xlsx":
        chunks = _xlsx(path, include_hidden, include_very_hidden)
    elif extension == ".xls":
        chunks = _xls(path, include_hidden, include_very_hidden)
    elif extension == ".docx":
        chunks = _docx(path)
    elif extension == ".csv":
        chunks = _csv(path)
    elif extension == ".txt":
        chunks = _text(path)
    else:
        text = extract_text(path)
        if extension in IMAGE_EXTENSIONS:
            chunks = [
                SourceChunk(value, "image", {"page_start": 1, "page_end": 1, "part": index})
                for index, value in enumerate(chunk_image_text(text), start=1)
            ]
        else:
            if extension == ".ppt":
                chunks = _split(text, "powerpoint", {
                    "slide_start": 1,
                    "slide_end": 1,
                    "shape_ids": ["legacy-text-stream"],
                    "shape_types": ["legacy_text"],
                    "speaker_notes_included": False,
                    "content_type": "legacy_text",
                })
            else:
                chunks = _split(text, "text", {
                    "line_start": 1,
                    "line_end": max(1, len(text.splitlines())),
                })
    if not chunks:
        raise DocumentParseError("No readable content was found in the uploaded file.")
    for chunk in chunks:
        validate_source_location(chunk.source_type, chunk.location)
    if monotonic() - started_at > settings.parser_timeout_seconds:
        raise DocumentParseError("Document parsing exceeded the configured time limit.")
    return chunks


def extract_source_metadata(
    path: Path,
    *,
    include_hidden: bool = True,
    include_very_hidden: bool = False,
) -> dict[str, object]:
    """Return workbook-level facts used for administrative status and QA."""
    if path.suffix.lower() not in {".xlsx", ".xls"}:
        return {}
    if path.suffix.lower() == ".xls":
        try:
            import xlrd

            workbook = xlrd.open_workbook(path, on_demand=True)
            try:
                sheets = workbook.sheets()
                supported = [
                    sheet for sheet in sheets
                    if (
                        int(getattr(sheet, "visibility", 0) or 0) == 0
                        or int(getattr(sheet, "visibility", 0) or 0) == 1
                        and include_hidden
                        or int(getattr(sheet, "visibility", 0) or 0) >= 2
                        and include_very_hidden
                    )
                ]
                return {
                    "source_type": "excel",
                    "sheet_count": len(sheets),
                    "sheet_names": [sheet.name for sheet in sheets],
                    "visible_sheet_count": sum(
                        not bool(getattr(sheet, "visibility", 0))
                        for sheet in sheets
                    ),
                    "processed_sheet_names": [sheet.name for sheet in supported],
                    "detected_tables": [],
                    "total_non_empty_rows": sum(
                        1
                        for sheet in supported
                        for row_index in range(sheet.nrows)
                        if any(
                            sheet.cell_value(row_index, column) not in ("", None)
                            for column in range(sheet.ncols)
                        )
                    ),
                }
            finally:
                workbook.release_resources()
        except Exception as error:
            raise DocumentParseError(
                "The Excel workbook metadata could not be read."
            ) from error
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=False, data_only=False)
        try:
            sheets = list(workbook.worksheets)
            supported = [
                sheet for sheet in sheets
                if (
                    sheet.sheet_state == "visible"
                    or sheet.sheet_state == "hidden" and include_hidden
                    or sheet.sheet_state == "veryHidden" and include_very_hidden
                )
            ]
            return {
                "source_type": "excel",
                "sheet_count": len(sheets),
                "sheet_names": [sheet.title for sheet in sheets],
                "visible_sheet_count": sum(
                    sheet.sheet_state == "visible" for sheet in sheets
                ),
                "processed_sheet_names": [sheet.title for sheet in supported],
                "detected_tables": [
                    {"sheet_name": sheet.title, "table_name": table.name}
                    for sheet in supported
                    for table in sheet.tables.values()
                ],
                "total_non_empty_rows": sum(
                    1
                    for sheet in supported
                    for row in sheet.iter_rows()
                    if any(cell.value is not None for cell in row)
                ),
            }
        finally:
            workbook.close()
    except Exception as error:
        raise DocumentParseError(
            "The Excel workbook metadata could not be read."
        ) from error
