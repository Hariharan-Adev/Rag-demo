"""Parser-registry tests for tabular, presentation, and OCR formats."""

from __future__ import annotations

import struct
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from app.services.document_loader import (
    DocumentParseError,
    ExcelParser,
    OcrParser,
    PARSER_REGISTRY,
    PowerPointParser,
    SUPPORTED_EXTENSIONS,
    _extract_legacy_ppt_text,
    extract_text,
)
from app.services.source_extraction import (
    extract_source_chunks,
    extract_source_metadata,
    validate_source_location,
)


class DocumentParserTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_registry_contains_every_required_extension(self):
        required = {".txt", ".pdf", ".docx", ".xlsx", ".xls", ".csv", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"}
        self.assertEqual(SUPPORTED_EXTENSIONS, required)
        self.assertIsInstance(PARSER_REGISTRY[".xlsx"], ExcelParser)
        self.assertIsInstance(PARSER_REGISTRY[".ppt"], PowerPointParser)
        self.assertIsInstance(PARSER_REGISTRY[".webp"], OcrParser)

    def test_csv_parser_extracts_rows(self):
        path = self.root / "employees.csv"
        path.write_text("Employee ID,Name,Department\n1001,John,HR\n1002,Alice,Finance", encoding="utf-8")
        text = extract_text(path)
        self.assertIn("CSV: employees", text)
        self.assertIn("Employee ID\tName\tDepartment", text)
        self.assertIn("1002\tAlice\tFinance", text)
        chunks = extract_source_chunks(path)
        self.assertIn("Employee ID: 1002", chunks[2].text)
        self.assertIn("Name: Alice", chunks[2].text)
        self.assertEqual(chunks[2].location["row_start"], 3)

    def test_xlsx_parser_extracts_workbook_sheets_and_cells(self):
        from openpyxl import Workbook

        path = self.root / "employees.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Employees"
        sheet.append(["Employee ID", "Name", "Department"])
        sheet.append([1001, "John", "HR"])
        workbook.create_sheet("Summary").append(["Total", 1])
        workbook.save(path)
        workbook.close()
        text = extract_text(path)
        self.assertIn("Workbook: employees", text)
        self.assertIn("Sheet: Employees", text)
        self.assertIn("1001\tJohn\tHR", text)
        self.assertIn("Sheet: Summary", text)

    def test_pptx_parser_extracts_slide_text_and_tables(self):
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "quarterly.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Quarterly Sales"
        slide.placeholders[1].text = "Revenue increased by 15%"
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(4), Inches(6), Inches(1)).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        presentation.save(path)
        text = extract_text(path)
        self.assertIn("Presentation: quarterly", text)
        self.assertIn("Slide 1", text)
        self.assertIn("Quarterly Sales", text)
        self.assertIn("Revenue increased by 15%", text)
        self.assertIn("Region\tRevenue", text)

    def test_pdf_chunks_retain_exact_page_numbers(self):
        path = self.root / "policy.pdf"
        path.write_bytes(b"%PDF-test")
        pages = [
            SimpleNamespace(extract_text=lambda: "First page"),
            SimpleNamespace(extract_text=lambda: "Second page"),
        ]
        with patch("pypdf.PdfReader", return_value=SimpleNamespace(pages=pages)):
            chunks = extract_source_chunks(path)
        self.assertEqual(
            [chunk.location["page_start"] for chunk in chunks],
            [1, 2],
        )
        self.assertTrue(all(chunk.source_type == "pdf" for chunk in chunks))

    def test_pptx_chunks_retain_slide_shape_and_table_row(self):
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "located.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Located title"
        table = slide.shapes.add_table(
            1, 2, Inches(1), Inches(4), Inches(6), Inches(1)
        ).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        presentation.save(path)
        chunks = extract_source_chunks(path)
        title = next(chunk for chunk in chunks if "Located title" in chunk.text)
        table_chunk = next(chunk for chunk in chunks if "Region" in chunk.text)
        self.assertEqual(title.location["slide_number"], 1)
        self.assertEqual(title.location["slide_start"], 1)
        self.assertEqual(title.location["slide_end"], 1)
        self.assertIn("shape_index", title.location)
        self.assertEqual(title.location["shape_ids"], ["shape-1"])
        self.assertFalse(title.location["speaker_notes_included"])
        self.assertEqual(table_chunk.location["row_start"], 1)
        self.assertEqual(table_chunk.location["content_type"], "table")

    def test_xlsx_chunks_retain_sheet_cell_range_table_and_formula(self):
        from openpyxl import Workbook
        from openpyxl.worksheet.table import Table

        path = self.root / "located.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Revenue"
        sheet.append(["Amount", "Tax"])
        sheet.append([100, "=A2*0.1"])
        sheet.add_table(Table(displayName="RevenueTable", ref="A1:B2"))
        workbook.save(path)
        workbook.close()
        chunks = extract_source_chunks(path)
        formula_row = next(chunk for chunk in chunks if "=A2*0.1" in chunk.text)
        self.assertIn("Sheet: Revenue", formula_row.text)
        self.assertIn("Amount: 100", formula_row.text)
        self.assertIn("Tax: =A2*0.1", formula_row.text)
        self.assertEqual(formula_row.location["sheet_name"], "Revenue")
        self.assertEqual(formula_row.location["cell_range"], "A2:B2")
        self.assertEqual(formula_row.location["table_name"], "RevenueTable")
        self.assertEqual(formula_row.location["formulas"], {"B2": "=A2*0.1"})
        self.assertEqual(formula_row.location["header_rows"], [1])
        self.assertEqual(formula_row.location["header_context"], ["Amount", "Tax"])
        self.assertFalse(formula_row.location["hidden_sheet"])
        metadata = extract_source_metadata(path)
        self.assertEqual(metadata["sheet_count"], 1)
        self.assertEqual(metadata["visible_sheet_count"], 1)
        self.assertEqual(metadata["total_non_empty_rows"], 2)
        self.assertEqual(
            metadata["detected_tables"],
            [{"sheet_name": "Revenue", "table_name": "RevenueTable"}],
        )

    def test_source_location_validation_rejects_incomplete_citations(self):
        with self.assertRaisesRegex(DocumentParseError, "missing page_end"):
            validate_source_location("pdf", {"page_start": 1})

    def test_legacy_ppt_text_atom_extraction(self):
        payload = "Legacy slide text".encode("utf-16-le")
        record = struct.pack("<HHI", 0, 4000, len(payload)) + payload
        self.assertEqual(_extract_legacy_ppt_text(record), ["Legacy slide text"])

    def test_ocr_parser_uses_all_image_frames(self):
        from PIL import Image

        path = self.root / "scan.png"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("pytesseract.image_to_string", return_value="Invoice total 42") as ocr,
            patch(
                "app.services.image_processor.image_parser.vision_is_configured",
                return_value=True,
            ),
            patch(
                "app.services.image_processor.image_parser.describe_image",
                return_value="A scanned invoice with a visible total.",
            ),
        ):
            text = extract_text(path)
        self.assertIn("Image: scan", text)
        self.assertIn("Invoice total 42", text)
        self.assertIn("A scanned invoice with a visible total.", text)
        ocr.assert_called_once()

    def test_missing_tesseract_returns_clear_error(self):
        import pytesseract
        from PIL import Image

        path = self.root / "scan.jpg"
        Image.new("RGB", (30, 20), "white").save(path)
        with (
            patch("pytesseract.image_to_string", side_effect=pytesseract.TesseractNotFoundError()),
            patch(
                "app.services.image_processor.image_parser.vision_is_configured",
                return_value=False,
            ),
        ):
            with self.assertRaisesRegex(DocumentParseError, "Tesseract service is not installed"):
                extract_text(path)

    def test_corrupt_excel_and_powerpoint_return_format_errors(self):
        xlsx = self.root / "broken.xlsx"
        pptx = self.root / "broken.pptx"
        xlsx.write_bytes(b"not a workbook")
        pptx.write_bytes(b"not a presentation")
        with self.assertRaisesRegex(DocumentParseError, "Excel workbook"):
            extract_text(xlsx)
        with self.assertRaisesRegex(DocumentParseError, "PowerPoint presentation"):
            extract_text(pptx)

    def test_unsupported_extension_is_rejected(self):
        path = self.root / "document.exe"
        path.write_bytes(b"unsafe")
        with self.assertRaisesRegex(DocumentParseError, "Unsupported document type"):
            extract_text(path)


if __name__ == "__main__":
    unittest.main()
